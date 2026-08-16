#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
read_any.py — DaShaAgent 统一多格式文件解析器。

设计参考 WorkBuddy 内核 Read 工具的三条核心范式:
  1) Token 预算而非字节截断 —— 超限时在"语义单元"边界切断, 并回传 next_offset 引导续读,
     而不是粗暴地砍掉一半汉字。
  2) 语义单元分页 —— PDF 按页 / PPTX 按幻灯片 / XLSX 按行 / EPUB 按章节 / 文本按行,
     offset+limit 让 agent 能在大文件里精确定位, 而不是"读不动就放弃"。
  3) 魔法字节嗅探优先, 扩展名兜底 —— 文件被改名也能正确解析。

输出: 单个 JSON 对象到 stdout (UTF-8)。永远输出 JSON, 即使失败, 便于上游稳定解析。

用法:
    python read_any.py <file> [--offset N] [--limit N] [--max-tokens N] [--mode text|meta]

信封格式:
{
  "ok": true/false,
  "kind": "pdf|docx|pptx|xlsx|epub|ipynb|html|rtf|eml|csv|json|zip|image|text|legacy|unknown",
  "path": "...",
  "meta": { ... },                # 格式相关元信息 (页数/表名/作者...)
  "unit": "page|slide|row|line|chapter|cell|entry|none",
  "total_units": 123,
  "offset": 1,                    # 本次返回的起始单元 (1-based)
  "returned_units": 20,
  "next_offset": 21,              # null 表示已读完
  "truncated": true/false,
  "est_tokens": 4321,
  "text": "...",
  "error": "CODE: message"        # ok=false 时存在
}
"""
import sys
import os
import io
import json
import re
import zipfile

# ---- Windows 控制台 UTF-8 兜底 ----
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

DEFAULT_MAX_TOKENS = 20000
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "webp", "tif", "tiff"}


# =========================================================================
# 工具函数
# =========================================================================

def est_tokens(s: str) -> int:
    """粗略 token 估算: ASCII 约 4 字符/token, CJK 约 1 字符/token。
    比 len(s)/4 对中文准确得多, 避免中文文档被严重低估后爆上下文。"""
    if not s:
        return 0
    ascii_n = 0
    wide_n = 0
    for ch in s:
        if ord(ch) < 128:
            ascii_n += 1
        else:
            wide_n += 1
    return int(ascii_n / 4) + wide_n


def decode_bytes(b: bytes) -> str:
    """无 chardet 时的编码探测: BOM -> utf-8 -> gb18030 -> latin-1。
    gb18030 是 gbk/gb2312 的超集, 覆盖国内绝大多数遗留文本。"""
    if b[:3] == b"\xef\xbb\xbf":
        return b[3:].decode("utf-8", errors="replace")
    if b[:2] in (b"\xff\xfe", b"\xfe\xff"):
        return b.decode("utf-16", errors="replace")
    for enc in ("utf-8", "gb18030", "big5", "latin-1"):
        try:
            return b.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return b.decode("utf-8", errors="replace")


def paginate(units, offset, limit, max_tokens, joiner="\n"):
    """在语义单元上做分页 + Token 预算双重约束。
    units: List[str], 每个元素是一个已格式化好的单元文本。
    返回 (text, returned_units, next_offset, truncated)"""
    total = len(units)
    start = max(1, int(offset or 1)) - 1
    if start >= total:
        return "", 0, None, False
    end_cap = total if not limit else min(total, start + int(limit))

    buf = []
    tok = 0
    i = start
    truncated = False
    while i < end_cap:
        u = units[i]
        t = est_tokens(u)
        # 至少放进一个单元, 否则空手而归
        if buf and tok + t > max_tokens:
            truncated = True
            break
        # 单个单元本身就超预算 -> 按字符硬切, 但保留提示
        if not buf and t > max_tokens:
            cut = max_tokens * 2  # 保守换算回字符
            u = u[:cut] + "\n…(该单元过大, 已截断)"
            buf.append(u)
            tok += est_tokens(u)
            i += 1
            truncated = True
            break
        buf.append(u)
        tok += t
        i += 1

    returned = i - start
    next_off = (i + 1) if i < total else None
    if next_off is None and limit and i < total:
        next_off = i + 1
    if i < total:
        truncated = True
    return joiner.join(buf), returned, next_off, truncated


def envelope(ok, kind, path, **kw):
    d = {"ok": ok, "kind": kind, "path": path}
    d.update(kw)
    return d


def emit(d):
    sys.stdout.write(json.dumps(d, ensure_ascii=False))
    sys.stdout.flush()


def fail(kind, path, code, msg, **kw):
    emit(envelope(False, kind, path, error="%s: %s" % (code, msg), text="", **kw))
    sys.exit(0)  # 始终 0 退出, 让上游读 JSON 判断成败, 避免 stderr/exitcode 双通道歧义


# =========================================================================
# 类型嗅探
# =========================================================================

def sniff(fp: str) -> str:
    try:
        with open(fp, "rb") as f:
            head = f.read(512)
    except Exception:
        head = b""

    ext = fp.lower().rsplit(".", 1)[-1] if "." in os.path.basename(fp) else ""

    # --- 魔法字节优先 ---
    if head[:4] == b"%PDF":
        return "pdf"
    if head[:8] == b"\x89PNG\r\n\x1a\n" or head[:3] == b"\xff\xd8\xff" \
            or head[:6] in (b"GIF87a", b"GIF89a") or head[:2] == b"BM" \
            or (head[:4] == b"RIFF" and head[8:12] == b"WEBP") \
            or head[:4] in (b"II*\x00", b"MM\x00*"):
        return "image"
    if head[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        # OLE2 复合文档: 老版 .doc/.xls/.ppt, 也可能是 .msg
        return "msg" if ext == "msg" else "legacy"
    if head[:5] == b"{\\rtf":
        return "rtf"
    if head[:2] == b"PK":
        # ZIP 容器族: docx/pptx/xlsx/epub/odt/jar/普通 zip
        try:
            with zipfile.ZipFile(fp) as z:
                names = set(z.namelist())
        except Exception:
            names = set()
        if "mimetype" in names or any(n.startswith("META-INF/container.xml") for n in names):
            return "epub"
        if any(n.startswith("word/") for n in names):
            return "docx"
        if any(n.startswith("ppt/") for n in names):
            return "pptx"
        if any(n.startswith("xl/") for n in names):
            return "xlsx"
        if ext in ("docx", "docm", "dotx", "dotm"):
            return "docx"
        if ext in ("pptx", "pptm", "ppsx", "potx"):
            return "pptx"
        if ext in ("xlsx", "xlsm", "xltx"):
            return "xlsx"
        if ext == "epub":
            return "epub"
        return "zip"

    # --- 扩展名兜底 ---
    if ext in IMAGE_EXTS:
        return "image"
    if ext == "ipynb":
        return "ipynb"
    if ext in ("html", "htm", "xhtml"):
        return "html"
    if ext in ("xml", "svg", "rss", "atom"):
        return "xml"
    if ext in ("eml", "mbox"):
        return "eml"
    if ext == "msg":
        return "msg"
    if ext == "rtf":
        return "rtf"
    if ext in ("csv", "tsv"):
        return "csv"
    if ext == "json":
        return "json"
    if ext in ("srt", "vtt", "ass"):
        return "subtitle"
    if ext == "pdf":
        return "pdf"
    if ext in ("doc", "xls", "ppt"):
        return "legacy"

    # --- 二进制探测: 有 NUL 字节基本不是文本 ---
    if b"\x00" in head:
        return "binary"
    return "text"


# =========================================================================
# 各格式解析器: 统一返回 (meta, unit_name, units:List[str])
# =========================================================================

def parse_pdf(fp):
    from pypdf import PdfReader
    reader = PdfReader(fp)
    n = len(reader.pages)
    meta = {"pages": n}
    try:
        info = reader.metadata or {}
        meta["title"] = str(info.get("/Title", "") or "")
        meta["author"] = str(info.get("/Author", "") or "")
    except Exception:
        pass

    units = []
    empty = 0
    for i, page in enumerate(reader.pages, 1):
        try:
            t = (page.extract_text() or "").strip()
        except Exception as e:
            t = "(第 %d 页解析失败: %s)" % (i, e)
        if not t:
            empty += 1
        units.append("----- 第 %d/%d 页 -----\n%s" % (i, n, t if t else "(空白页或纯图像)"))

    # 文字层缺失比例过高 -> pdfplumber 兜底重抽
    if n and empty / n > 0.6:
        try:
            import pdfplumber
            alt = []
            with pdfplumber.open(fp) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    t = (page.extract_text() or "").strip()
                    alt.append("----- 第 %d/%d 页 -----\n%s" % (i, n, t if t else "(空白页或纯图像)"))
            joined_alt = sum(len(x) for x in alt)
            joined_old = sum(len(x) for x in units)
            if joined_alt > joined_old * 1.2:
                units = alt
                meta["extractor"] = "pdfplumber"
        except Exception:
            pass
        meta["hint"] = ("该 PDF 大部分页面无文字层, 可能是扫描件。"
                        "如需正文, 请用 run_code 调 PyMuPDF 渲染成图后 OCR。")
    meta.setdefault("extractor", "pypdf")
    return meta, "page", units


def parse_docx(fp):
    from docx import Document
    doc = Document(fp)
    units = []
    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        style = (p.style.name or "") if p.style else ""
        if style.startswith("Heading") or style.startswith("标题"):
            lvl = re.findall(r"\d+", style)
            hashes = "#" * (int(lvl[0]) if lvl else 1)
            units.append("%s %s" % (hashes, text))
        else:
            units.append(text)
    for ti, table in enumerate(doc.tables, 1):
        rows = ["[表格 %d]" % ti]
        for row in table.rows:
            rows.append(" | ".join((c.text or "").strip() for c in row.cells))
        units.append("\n".join(rows))
    meta = {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
    try:
        cp = doc.core_properties
        meta["title"] = cp.title or ""
        meta["author"] = cp.author or ""
    except Exception:
        pass
    return meta, "block", units


def parse_pptx(fp):
    from pptx import Presentation
    prs = Presentation(fp)
    slides = list(prs.slides)
    units = []
    for i, slide in enumerate(slides, 1):
        parts = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for r in shape.table.rows:
                        parts.append(" | ".join((c.text or "").strip() for c in r.cells))
            except Exception:
                continue
        # 备注页
        try:
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    parts.append("[备注] " + nt)
        except Exception:
            pass
        units.append("----- 幻灯片 %d/%d -----\n%s" % (i, len(slides), "\n".join(parts) or "(无文本)"))
    return {"slides": len(slides)}, "slide", units


def parse_xlsx(fp):
    from openpyxl import load_workbook
    wb = load_workbook(fp, data_only=True, read_only=True)
    units = []
    sheet_info = []
    for sn in wb.sheetnames:
        ws = wb[sn]
        mr, mc = ws.max_row or 0, ws.max_column or 0
        sheet_info.append({"name": sn, "rows": mr, "cols": mc})
        units.append("=== 工作表: %s (%d行 x %d列) ===" % (sn, mr, mc))
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            line = " | ".join(vals).rstrip(" |")
            if line.strip():
                units.append(line)
    wb.close()
    return {"sheets": sheet_info}, "row", units


def parse_epub(fp):
    """不依赖 ebooklib: epub = ZIP + OPF 清单 + XHTML 章节。
    按 OPF spine 顺序还原真实阅读顺序, 而不是按 zip 内文件名瞎排。"""
    from bs4 import BeautifulSoup
    units = []
    meta = {}
    with zipfile.ZipFile(fp) as z:
        names = z.namelist()
        # 1. 定位 OPF
        opf_path = None
        try:
            container = z.read("META-INF/container.xml").decode("utf-8", "replace")
            m = re.search(r'full-path="([^"]+)"', container)
            if m:
                opf_path = m.group(1)
        except Exception:
            pass
        if not opf_path:
            cand = [n for n in names if n.lower().endswith(".opf")]
            opf_path = cand[0] if cand else None

        spine_files = []
        if opf_path:
            base = os.path.dirname(opf_path)
            opf = z.read(opf_path).decode("utf-8", "replace")
            osoup = BeautifulSoup(opf, "xml")
            for tag in ("title", "creator", "language", "publisher", "date"):
                el = osoup.find(tag)
                if el and el.get_text(strip=True):
                    meta[tag] = el.get_text(strip=True)
            manifest = {}
            for item in osoup.find_all("item"):
                iid = item.get("id")
                href = item.get("href")
                if iid and href:
                    manifest[iid] = os.path.join(base, href).replace("\\", "/") if base else href
            for ref in osoup.find_all("itemref"):
                idref = ref.get("idref")
                if idref and idref in manifest:
                    spine_files.append(manifest[idref])
        if not spine_files:
            spine_files = sorted(n for n in names
                                 if n.lower().endswith((".xhtml", ".html", ".htm")))

        meta["chapters"] = len(spine_files)
        for idx, sf in enumerate(spine_files, 1):
            try:
                raw = z.read(sf)
            except KeyError:
                continue
            soup = BeautifulSoup(decode_bytes(raw), "lxml")
            for bad in soup(["script", "style"]):
                bad.decompose()
            # 章节标题: 正文里的 h1~h3 优先, <head><title> 只作兜底
            h = soup.find(["h1", "h2", "h3", "h4"])
            title = h.get_text(strip=True) if h else None
            if not title:
                t = soup.find("title")
                title = t.get_text(strip=True) if t else os.path.basename(sf)
            # 只取 <body>, 否则 <head><title> 会被 get_text 一起抓进正文造成重复
            root = soup.find("body") or soup
            for bad in root.find_all("title"):
                bad.decompose()
            body = root.get_text("\n", strip=True)
            body = re.sub(r"\n{3,}", "\n\n", body)
            units.append("----- 第 %d 章 · %s -----\n%s" % (idx, title, body))
    return meta, "chapter", units


def parse_ipynb(fp):
    """不依赖 nbformat: ipynb 就是 JSON。"""
    with open(fp, "rb") as f:
        nb = json.loads(decode_bytes(f.read()))
    cells = nb.get("cells", [])
    lang = (nb.get("metadata", {}).get("kernelspec", {}) or {}).get("language", "python")
    units = []
    for i, c in enumerate(cells, 1):
        ctype = c.get("cell_type", "?")
        src = c.get("source", [])
        body = "".join(src) if isinstance(src, list) else str(src)
        head = "----- Cell %d [%s] -----" % (i, ctype)
        chunk = [head, body.rstrip()]
        # 只保留文本型输出, 图片/二进制输出跳过 (避免 base64 灌爆上下文)
        outs = c.get("outputs", []) or []
        otexts = []
        for o in outs:
            if o.get("output_type") == "stream":
                t = o.get("text", [])
                otexts.append("".join(t) if isinstance(t, list) else str(t))
            elif o.get("output_type") in ("execute_result", "display_data"):
                d = o.get("data", {})
                if "text/plain" in d:
                    t = d["text/plain"]
                    otexts.append("".join(t) if isinstance(t, list) else str(t))
                elif any(k.startswith("image/") for k in d):
                    otexts.append("(图像输出, 已省略)")
            elif o.get("output_type") == "error":
                otexts.append("ERROR %s: %s" % (o.get("ename"), o.get("evalue")))
        if otexts:
            chunk.append("  >> 输出:\n" + "\n".join(otexts)[:2000])
        units.append("\n".join(chunk))
    return {"cells": len(cells), "language": lang}, "cell", units


def parse_html(fp, is_xml=False):
    from bs4 import BeautifulSoup
    with open(fp, "rb") as f:
        raw = f.read()
    soup = BeautifulSoup(decode_bytes(raw), "xml" if is_xml else "lxml")
    meta = {}
    if not is_xml:
        if soup.title and soup.title.string:
            meta["title"] = soup.title.string.strip()
        for bad in soup(["script", "style", "noscript"]):
            bad.decompose()
        # 提取链接清单, 对 agent 抓取后续目标很有用
        links = []
        for a in soup.find_all("a", href=True)[:80]:
            txt = a.get_text(strip=True)
            if txt:
                links.append("%s -> %s" % (txt[:60], a["href"]))
        if links:
            meta["links_sample"] = links[:30]
    text = soup.get_text("\n", strip=True)
    text = re.sub(r"\n{3,}", "\n\n", text)
    units = text.split("\n")
    return meta, "line", units


def parse_rtf(fp):
    """不依赖 striprtf: 手写 RTF 控制字剥离。
    处理 \\uN 转义(中文)、\\'hh 十六进制、忽略图片/字体表等 destination 组。"""
    with open(fp, "rb") as f:
        raw = f.read().decode("latin-1", "replace")

    out = []
    i, n = 0, len(raw)
    depth = 0
    # skip_from: 一旦某个 destination 组被判定为"丢弃", 该组及其所有嵌套子组
    # 都必须丢弃。只记住起始深度即可, 比用 set 逐层记更正确 —— 之前用 set 时
    # {\fonttbl{\f0 Arial;}} 的内层 {\f0 ...} 深度不同, 会漏出 "Arial;"。
    skip_from = None
    SKIP_DESTS = ("fonttbl", "colortbl", "stylesheet", "info", "pict",
                  "object", "themedata", "colorschememapping", "latentstyles",
                  "datastore", "generator", "listtable", "rsidtbl", "xmlnstbl",
                  "filetbl", "revtbl", "upr", "header", "footer", "footnote",
                  "mmathPr", "wgrffmtfilter", "shppict", "nonshppict")

    def live():
        return skip_from is None

    while i < n:
        ch = raw[i]
        if ch == "{":
            depth += 1
            i += 1
            continue
        if ch == "}":
            depth -= 1
            if skip_from is not None and depth < skip_from:
                skip_from = None
            i += 1
            continue
        if ch == "\\":
            # 转义字符 \\ \{ \}
            if i + 1 < n and raw[i + 1] in "\\{}":
                if live():
                    out.append(raw[i + 1])
                i += 2
                continue
            # \'hh 十六进制字节 (ANSI 代码页, 中文常见)
            if raw[i + 1:i + 2] == "'":
                hexs = raw[i + 2:i + 4]
                try:
                    b = bytes([int(hexs, 16)])
                    if live():
                        out.append(b.decode("gb18030", "replace"))
                except Exception:
                    pass
                i += 4
                continue
            m = re.match(r"\\([a-zA-Z]+)(-?\d+)?[ ]?", raw[i:])
            if m:
                word, param = m.group(1), m.group(2)
                # \uN Unicode 码点
                if word == "u" and param is not None:
                    if live():
                        cp = int(param)
                        if cp < 0:
                            cp += 65536
                        try:
                            out.append(chr(cp))
                        except Exception:
                            pass
                    i += m.end()
                    if i < n and raw[i] == "?":   # 跳过后随的 ANSI 替代字符
                        i += 1
                    continue
                if word in SKIP_DESTS:
                    if skip_from is None:
                        skip_from = depth
                elif live():
                    if word in ("par", "line", "sect", "page"):
                        out.append("\n")
                    elif word == "tab":
                        out.append("\t")
                    elif word in ("emdash", "endash"):
                        out.append("—")
                    elif word in ("lquote", "rquote"):
                        out.append("'")
                    elif word in ("ldblquote", "rdblquote"):
                        out.append('"')
                i += m.end()
                continue
            i += 1
            continue
        # \* 开头的未知 destination 也应丢弃
        if ch == "*" and raw[i - 1:i] == "\\":
            if skip_from is None:
                skip_from = depth
            i += 1
            continue
        if live() and ch not in "\r\n":
            out.append(ch)
        i += 1

    text = "".join(out)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    units = [l for l in text.split("\n")]
    return {"note": "RTF 由内置解析器剥离控制字, 复杂排版可能丢失"}, "line", units


def parse_eml(fp):
    """标准库 email 模块, 支持 multipart / 附件清单 / 编码头解码。"""
    from email import policy
    from email.parser import BytesParser
    with open(fp, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    meta = {}
    for h in ("From", "To", "Cc", "Subject", "Date", "Message-ID"):
        v = msg.get(h)
        if v:
            meta[h.lower()] = str(v)

    units = []
    for h in ("From", "To", "Cc", "Subject", "Date"):
        if msg.get(h):
            units.append("%s: %s" % (h, msg.get(h)))
    units.append("-" * 40)

    attachments = []
    body_text = ""
    if msg.is_multipart():
        for part in msg.walk():
            disp = part.get_content_disposition()
            ctype = part.get_content_type()
            if disp == "attachment":
                attachments.append("%s (%s)" % (part.get_filename() or "unnamed", ctype))
                continue
            if ctype == "text/plain" and not body_text:
                body_text = part.get_content()
            elif ctype == "text/html" and not body_text:
                try:
                    from bs4 import BeautifulSoup
                    body_text = BeautifulSoup(part.get_content(), "lxml").get_text("\n", strip=True)
                except Exception:
                    body_text = part.get_content()
    else:
        try:
            body_text = msg.get_content()
        except Exception:
            body_text = decode_bytes(msg.get_payload(decode=True) or b"")

    units.extend(str(body_text).split("\n"))
    if attachments:
        meta["attachments"] = attachments
        units.append("\n----- 附件 (%d) -----" % len(attachments))
        units.extend(attachments)
    return meta, "line", units


def parse_zip(fp):
    with zipfile.ZipFile(fp) as z:
        infos = z.infolist()
    units = ["=== ZIP 归档 · 共 %d 个条目 ===" % len(infos)]
    total = 0
    for info in infos:
        total += info.file_size
        units.append("%-60s %10d bytes" % (info.filename[:60], info.file_size))
    return {"entries": len(infos), "uncompressed_bytes": total,
            "hint": "这是压缩包清单。需要读取内部文件请用 run_code 解压后再 fs_read。"}, "entry", units


MIME_BY_MAGIC = [
    (b"\x89PNG\r\n\x1a\n", "image/png"),
    (b"\xff\xd8\xff", "image/jpeg"),
    (b"GIF87a", "image/gif"),
    (b"GIF89a", "image/gif"),
    (b"BM", "image/bmp"),
]


def _image_mime(head: bytes, fp: str) -> str:
    for magic, mime in MIME_BY_MAGIC:
        if head.startswith(magic):
            return mime
    if head[:4] == b"RIFF" and head[8:12] == b"WEBP":
        return "image/webp"
    if head[:4] in (b"II*\x00", b"MM\x00*"):
        return "image/tiff"
    ext = fp.lower().rsplit(".", 1)[-1]
    return {"jpg": "image/jpeg", "jpeg": "image/jpeg"}.get(ext, "image/" + (ext or "png"))


def encode_image_b64(fp, max_side=1568, max_bytes=3_500_000):
    """把图片编码为 base64, 供多模态模型直读。

    这是 WorkBuddy 内核 Read 工具处理图片的首选路径 —— 保真度远高于 OCR。
    两道瘦身闸门:
      1) 长边超过 max_side 等比缩放 (视觉模型普遍在 ~1.5k 长边后收益递减, 但 token 线性涨);
      2) 编码后仍超 max_bytes 则逐档降 JPEG 质量。
    PIL 不可用时原样回传 (小图无妨, 大图给出告警)。
    """
    import base64
    with open(fp, "rb") as f:
        raw = f.read()
    head = raw[:16]
    mime = _image_mime(head, fp)
    meta = {"orig_bytes": len(raw), "mime": mime}

    try:
        from PIL import Image
    except Exception:
        meta["resize"] = "PIL 不可用, 原样传输"
        if len(raw) > max_bytes:
            meta["warn"] = "图片过大且无法压缩, 可能超出模型输入上限"
        return base64.b64encode(raw).decode("ascii"), mime, meta

    try:
        with Image.open(fp) as im:
            meta["width"], meta["height"] = im.size
            need_resize = max(im.size) > max_side
            if not need_resize and len(raw) <= max_bytes:
                return base64.b64encode(raw).decode("ascii"), mime, meta

            im = im.convert("RGB") if im.mode not in ("RGB", "L") else im
            if need_resize:
                ratio = max_side / float(max(im.size))
                new_size = (max(1, int(im.width * ratio)), max(1, int(im.height * ratio)))
                im = im.resize(new_size, Image.LANCZOS)
                meta["resized_to"] = list(new_size)

            for q in (85, 70, 55, 40):
                buf = io.BytesIO()
                im.save(buf, format="JPEG", quality=q, optimize=True)
                data = buf.getvalue()
                if len(data) <= max_bytes:
                    meta["jpeg_quality"] = q
                    meta["encoded_bytes"] = len(data)
                    return base64.b64encode(data).decode("ascii"), "image/jpeg", meta
            meta["jpeg_quality"] = 40
            meta["encoded_bytes"] = len(data)
            meta["warn"] = "已最大压缩仍偏大"
            return base64.b64encode(data).decode("ascii"), "image/jpeg", meta
    except Exception as e:
        meta["encode_error"] = str(e)
        return base64.b64encode(raw).decode("ascii"), mime, meta


def parse_image(fp, allow_ocr=True):
    """图片: 优先返回元信息 + OCR 文本。
    注意: WorkBuddy 内核的首选是把图片 base64 直接喂给多模态模型,
    OCR 只是模型不支持图片时的兜底。这里返回 ocr_available 标志,
    由上游 fs_read 决定走哪条路。"""
    meta = {}
    try:
        from PIL import Image
        with Image.open(fp) as im:
            meta["width"], meta["height"] = im.size
            meta["format"] = im.format
            meta["mode"] = im.mode
    except Exception as e:
        meta["pil_error"] = str(e)

    if not allow_ocr:
        return meta, "none", ["(图片文件, 未执行 OCR)"]

    text = ""
    try:
        import pytesseract
        from PIL import Image
        # Windows 常见安装路径兜底
        for p in (r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                  r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"):
            if os.path.exists(p):
                pytesseract.pytesseract.tesseract_cmd = p
                break
        with Image.open(fp) as im:
            if im.mode not in ("RGB", "L"):
                im = im.convert("RGB")
            try:
                text = pytesseract.image_to_string(im, lang="chi_sim+eng")
            except Exception:
                text = pytesseract.image_to_string(im)
        meta["ocr"] = "tesseract"
    except Exception as e:
        meta["ocr_error"] = str(e)
        meta["hint"] = "OCR 不可用。若当前模型支持多模态, 应改用图片直读通道。"

    text = (text or "").strip()
    if not text:
        return meta, "none", ["(图片中未识别到文字)"]
    return meta, "line", text.split("\n")


def parse_csv(fp):
    import csv as csvmod
    with open(fp, "rb") as f:
        raw = f.read()
    txt = decode_bytes(raw)
    sample = txt[:4096]
    try:
        dialect = csvmod.Sniffer().sniff(sample, delimiters=",;\t|")
        delim = dialect.delimiter
    except Exception:
        delim = "\t" if fp.lower().endswith(".tsv") else ","
    rows = list(csvmod.reader(io.StringIO(txt), delimiter=delim))
    units = [" | ".join(r) for r in rows]
    meta = {"rows": len(rows), "delimiter": delim,
            "columns": len(rows[0]) if rows else 0,
            "header": rows[0] if rows else []}
    return meta, "row", units


def parse_text(fp):
    with open(fp, "rb") as f:
        raw = f.read()
    txt = decode_bytes(raw)
    units = txt.split("\n")
    return {"bytes": len(raw), "lines": len(units)}, "line", units


# =========================================================================
# 主流程
# =========================================================================

DISPATCH = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "epub": parse_epub,
    "ipynb": parse_ipynb,
    "rtf": parse_rtf,
    "eml": parse_eml,
    "zip": parse_zip,
    "csv": parse_csv,
    "text": parse_text,
    "json": parse_text,
    "subtitle": parse_text,
}


def main():
    args = sys.argv[1:]
    if not args:
        fail("unknown", "", "USAGE", "read_any.py <file> [--offset N] [--limit N] [--max-tokens N]")

    src = args[0]
    offset, limit, max_tokens = 1, 0, DEFAULT_MAX_TOKENS
    mode = "text"
    no_ocr = False
    i = 1
    while i < len(args):
        a = args[i]
        if a == "--offset" and i + 1 < len(args):
            offset = int(args[i + 1]); i += 2
        elif a == "--limit" and i + 1 < len(args):
            limit = int(args[i + 1]); i += 2
        elif a == "--max-tokens" and i + 1 < len(args):
            max_tokens = int(args[i + 1]); i += 2
        elif a == "--mode" and i + 1 < len(args):
            mode = args[i + 1]; i += 2
        elif a == "--no-ocr":
            no_ocr = True; i += 1
        elif a == "--b64":
            mode = "b64"; i += 1
        else:
            i += 1

    if not os.path.exists(src):
        fail("unknown", src, "FILE_NOT_FOUND", src)
    if os.path.isdir(src):
        fail("unknown", src, "IS_DIRECTORY", "这是目录, 请用 fs_list 列目录")

    kind = sniff(src)
    size = os.path.getsize(src)

    # --b64: 图片 base64 直读通道 (多模态模型优先路径, 不做 OCR)
    if mode == "b64":
        if kind != "image":
            fail(kind, src, "NOT_AN_IMAGE", "--b64 仅用于图片文件", meta={"bytes": size})
        try:
            b64, mime, imeta = encode_image_b64(src)
        except Exception as e:
            fail("image", src, "ENCODE_ERROR", "%s: %s" % (type(e).__name__, e), meta={"bytes": size})
        imeta["bytes"] = size
        emit(envelope(True, "image", src, meta=imeta, unit="none",
                      total_units=1, offset=1, returned_units=1, next_offset=None,
                      truncated=False, est_tokens=0, text="", b64=b64, mime=mime))
        return

    if kind == "legacy":
        fail("legacy", src, "LEGACY_OLE_FORMAT",
             "检测到旧版 Office 复合文档(.doc/.xls/.ppt)。请先另存为 .docx/.xlsx/.pptx, "
             "或用 run_code 调用 LibreOffice/antiword 转换后再读。",
             meta={"bytes": size})
    if kind == "msg":
        fail("msg", src, "MSG_UNSUPPORTED",
             "Outlook .msg 需要 extract_msg 库(当前未安装)。可先在 Outlook 中另存为 .eml 再读取。",
             meta={"bytes": size})
    if kind == "binary":
        fail("binary", src, "BINARY_FILE",
             "这是二进制文件, 无文本内容可抽取。若确需查看请用 run_code 做十六进制转储。",
             meta={"bytes": size})

    try:
        if kind == "image":
            meta, unit, units = parse_image(src, allow_ocr=not no_ocr)
        elif kind in ("html", "xml"):
            meta, unit, units = parse_html(src, is_xml=(kind == "xml"))
        else:
            fn = DISPATCH.get(kind, parse_text)
            meta, unit, units = fn(src)
    except ImportError as e:
        fail(kind, src, "MISSING_DEPENDENCY", str(e), meta={"bytes": size})
    except Exception as e:
        fail(kind, src, "PARSE_ERROR", "%s: %s" % (type(e).__name__, e), meta={"bytes": size})

    meta = meta or {}
    meta["bytes"] = size

    if mode == "meta":
        emit(envelope(True, kind, src, meta=meta, unit=unit,
                      total_units=len(units), offset=0, returned_units=0,
                      next_offset=1, truncated=False, est_tokens=0, text=""))
        return

    text, returned, next_off, truncated = paginate(units, offset, limit, max_tokens)

    emit(envelope(True, kind, src,
                  meta=meta,
                  unit=unit,
                  total_units=len(units),
                  offset=offset,
                  returned_units=returned,
                  next_offset=next_off,
                  truncated=truncated,
                  est_tokens=est_tokens(text),
                  text=text))


if __name__ == "__main__":
    main()
