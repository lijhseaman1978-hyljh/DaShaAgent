#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
read_office.py — 把 .pdf / .docx / .pptx 抽取为可读纯文本。
参考 WorkBuddy 内置 office 技能「按内容品类(doc/sheet/slide)归类」的思路，
这里用「魔法字节(magic bytes)嗅探」识别文件真实类型，扩展名仅作兜底，
比单纯看后缀更鲁棒（避免用户把 .pdf 改名成 .txt 之类的情况）。

用法:
    python read_office.py <file.pdf|.docx|.pptx>
输出到 stdout；失败信息以 "FAIL: " 前缀打印到 stdout 并以非 0 退出。
"""
import sys
import os

MAX_CHARS = 120_000  # 与 harness 侧 Observation 截断保持一致量级的兜底


def sniff_category(fp: str) -> str:
    """按文件头部魔法字节判断品类，扩展名兜底。"""
    try:
        with open(fp, "rb") as f:
            head = f.read(8)
    except Exception:
        head = b""
    if head[:4] == b"%PDF":
        return "pdf"
    if head[:2] == b"PK":  # ZIP 容器：docx / pptx / xlsx / odt ...
        # 进一步用扩展名区分 docx vs pptx（二者都是 ZIP，靠内部 Content_Types 也可，但扩展名足够）
        ext = fp.lower().rsplit(".", 1)[-1] if "." in fp else ""
        if ext in ("docx", "docm", "dotx", "dotm"):
            return "docx"
        if ext in ("pptx", "pptm", "ppsx", "potx"):
            return "pptx"
        # 无扩展名或未知 ZIP：尝试用内部 [Content_Types].xml 判断
        try:
            import zipfile
            with zipfile.ZipFile(fp) as z:
                names = z.namelist()
            if any(n.endswith("presentation.xml") for n in names):
                return "pptx"
            if any(n.endswith("document.xml") for n in names):
                return "docx"
        except Exception:
            pass
        return "docx" if ext.startswith("doc") else ("pptx" if ext.startswith("ppt") else "unknown")
    if head[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # OLE2：老版 .doc/.xls/.ppt
        return "legacy"
    # 兜底：看扩展名
    ext = fp.lower().rsplit(".", 1)[-1] if "." in fp else ""
    if ext == "pdf":
        return "pdf"
    if ext in ("docx", "docm", "dotx", "dotm"):
        return "docx"
    if ext in ("pptx", "pptm", "ppsx", "potx"):
        return "pptx"
    return "unknown"


def read_pdf(fp: str) -> str:
    out = []
    try:
        from pypdf import PdfReader
        reader = PdfReader(fp)
        out.append("=== PDF · 共 %d 页 ===" % len(reader.pages))
        for i, page in enumerate(reader.pages, 1):
            txt = page.extract_text() or ""
            out.append("\n----- 第 %d 页 -----\n%s" % (i, txt.strip()))
        joined = "\n".join(out)
        # 若 pypdf 抽不到（扫描件/图片型 PDF），用 pdfplumber 兜底一页页重抽
        if len(joined.strip()) < 40:
            try:
                import pdfplumber
                out2 = ["=== PDF( pdfplumber 兜底 ) · 共 ? 页 ==="]
                with pdfplumber.open(fp) as pdf:
                    for i, page in enumerate(pdf.pages, 1):
                        t = page.extract_text() or ""
                        out2.append("\n----- 第 %d 页 -----\n%s" % (i, t.strip()))
                joined = "\n".join(out2)
            except Exception:
                pass
        return joined
    except Exception as e:
        return "FAIL: PDF_PARSE_ERROR: " + str(e)


def read_docx(fp: str) -> str:
    try:
        from docx import Document
        doc = Document(fp)
        out = ["=== Word 文档 (.docx) ==="]
        for p in doc.paragraphs:
            style = (p.style.name or "") if p.style else ""
            text = p.text or ""
            if not text.strip():
                continue
            # 标题加标记，提升可读性
            if style.startswith("Heading") or style.startswith("标题"):
                out.append("\n# " + text)
            else:
                out.append(text)
        # 表格
        if doc.tables:
            out.append("\n----- 表格 -----")
            for ti, table in enumerate(doc.tables, 1):
                out.append("[表格 %d]" % ti)
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        return "\n".join(out)
    except Exception as e:
        return "FAIL: DOCX_PARSE_ERROR: " + str(e)


def read_pptx(fp: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(fp)
        out = ["=== PowerPoint 演示 (.pptx) · 共 %d 张幻灯片 ===" % len(prs.slides)]
        for i, slide in enumerate(prs.slides, 1):
            out.append("\n----- 幻灯片 %d -----" % i)
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
                if shape.has_table:
                    for r in shape.table.rows:
                        cells = [c.text.strip() for c in r.cells]
                        texts.append(" | ".join(cells))
            out.append("\n".join(texts) if texts else "(无文本)")
        return "\n".join(out)
    except Exception as e:
        return "FAIL: PPTX_PARSE_ERROR: " + str(e)


def main():
    if len(sys.argv) < 2:
        print("FAIL: usage: read_office.py <file.pdf|.docx|.pptx>")
        sys.exit(2)
    src = sys.argv[1]
    if not os.path.exists(src):
        print("FAIL: FILE_NOT_FOUND: %s" % src)
        sys.exit(3)

    cat = sniff_category(src)
    if cat == "pdf":
        text = read_pdf(src)
    elif cat == "docx":
        text = read_docx(src)
    elif cat == "pptx":
        text = read_pptx(src)
    elif cat == "legacy":
        text = ("FAIL: LEGACY_FORMAT: 检测到旧版 Office 格式(.doc/.ppt/.xls，OLE 复合文档)。\n"
                "当前解析器仅支持新版 OpenXML(.docx/.pptx/.xlsx)。请先用 Word/WPS 另存为 .docx/.pptx 后再读取。")
    else:
        text = "FAIL: UNSUPPORTED: 无法识别文件类型（非 PDF/DOCX/PPTX）。如需读取，请用 fs_read 文本模式或自写解析脚本后用 run_code 运行。"

    if text.startswith("FAIL:"):
        sys.stdout.write(text)
        sys.exit(5)

    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + "\n…(内容过长，仅显示前 %d 字符)" % MAX_CHARS
    sys.stdout.write(text)


if __name__ == "__main__":
    main()
